# %%
import base64
import math
import os


# NEO4J setup
os.environ['NEO4J_LOGIN'] = "neo4j"
os.environ['NEO4J_URL'] = "bolt://192.168.0.17:7687"
os.environ['NEO4J_PASSWORD'] = "1234554321"


from time import sleep
from pptx import Presentation
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches
from pptx.shapes.autoshape import Shape
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

from app.graph.connector2 import NeoSession, neodbsession
from credentials import *

from selenium import webdriver
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.common.by import By
from selenium.webdriver.support.wait import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
import io

from pptx import Presentation


class CypherPath:
    INSIGHT_FROM_CLOSEST_NEAR_INTENTION = '''
        match p2=((n:H_UPDATE {text:$insight, label:'insight'})-[:LEADS_TO]->(i)-[:INSIGHT*0..20]-(j)-[:LEADS_TO]-(n1:H_UPDATE {label:'intention'}))
        match p1=(
            (n1)-[:LEADS_TO]-()-[:FEEDBACK]-(c)-[:UPDATE*0..20]-(m:C_STATE)-[:FEEDBACK]-()-[:LEADS_TO]-(n)
            )
        WITH nodes(p1) as px order by length(p1) limit 1
        UNWIND px as mainPath
        CALL{
            with mainPath
            optional match (mainPath)-[:LEADS_TO]-(interactions:C_UPDATE)
            return interactions LIMIT 1
        }
        return mainPath, interactions
        '''
    INSIGHT_FROM_LONGEST_NEAR_INTENTION = '''
        match p2=((n:H_UPDATE {text:$insight, label:'insight'})-[:LEADS_TO]->(i)<-[:INSIGHT*0..20]-(j)-[:LEADS_TO]-(n1:H_UPDATE {label:'intention'}))
        match p1=(
            (n1)-[:LEADS_TO]-()-[:FEEDBACK]-(c)-[:UPDATE*0..20]-(m:C_STATE)-[:FEEDBACK]-()-[:LEADS_TO]-(n)
            )
        WITH nodes(p1) as px order by length(p1) desc limit 1
        UNWIND px as mainPath
        optional match (mainPath)-[:LEADS_TO]-(interactions:C_UPDATE)
        return mainPath, interactions
        '''
    ONLY_INSIGHT = '''
        match p1=((n:H_UPDATE {text:$insight, label:'insight'}))
        return p1 limit 1
        '''
    ALL_INSIGHTS_FROM_INTENTION = '''MATCH intention_path=
    ((n:H_UPDATE {text:$intention, label:'intention', tool:'qol'})
    -[:LEADS_TO]-(i)-[:INSIGHT*0..20]-(j)-[:LEADS_TO]
    -(n1:H_UPDATE {label:'insight', tool:'qol'}))
    WITH nodes(intention_path) as px
    UNWIND px as mainPath
    return mainPath, count(mainPath)'''
    INSIGHT_INTENTION = '''
        match p1=((n1:H_UPDATE {text:$intention, label:'intention'})-[:LEADS_TO]-()-[:FEEDBACK]-(c)-[:UPDATE*0..20]-(m:C_STATE)-[:FEEDBACK]-()-[:LEADS_TO]-(n:H_UPDATE {text:$insight, label:'insight'}))
        WITH nodes(p1) as px order by length(p1) limit 1
        UNWIND px as mainPath
        CALL{
            with mainPath
            optional match (mainPath)-[:LEADS_TO]-(interactions:C_UPDATE)
            return interactions LIMIT 1
        }
        return mainPath, interactions
        '''

    def run(self, query, **args):
        db = neodbsession.engine.begin()
        path = db.run(query, args).data()
        # print(path)
        # nodes = path[0]['p1'].nodes[::-1]
        neodbsession.engine.rollback(db)
        return path


# def get_path(start):
#     db = neodbsession.engine.begin()
#     path = db.run(f'''
#     match p=((n:H_UPDATE {{text:'{start}'}})<-[:DOES]-()<-[:INSIGHT*0..20]-()-[:LEADS_TO]-(n1:H_UPDATE {{label:'intention'}}))
#     match p1=(
#         (n)-[:LEADS_TO]-()-[:FEEDBACK]-(m:C_STATE)
#         -[:UPDATE*0..20]-(c)-[:FEEDBACK]-()-[:LEADS_TO]-(n1)
#         )
#     return p1
#     ''').data()
#     print(path)
#     nodes = path[0]['p1'].nodes[::-1]
#     neodbsession.engine.rollback(db)
#     return nodes


class Chrome:
    def __init__(self, wait=5):
        self.driver = webdriver.Chrome('../local_storage/chromedriver.exe')
        self.url = 'http://localhost:3000/'
        # self.qolurl = 'http://localhost:3000/'
        # setup to fetch proper website screenshot
        self.qolurl = 'https://qol.vav.aknakos.com/'
        self.oceanurl = 'https://ocean.aknakos.com/'
        self.wait = wait

    def setup_pic_wbt(self):
        self.driver.get(self.qolurl)
        element = WebDriverWait(self.driver, 10).until(
            EC.presence_of_element_located((By.ID, 'disable-tutorial'))
        )
        element.click()
        element = WebDriverWait(self.driver, 10).until(
            EC.presence_of_element_located((By.TAG_NAME, "path"))
        )

    def setup_pic_ocean(self):
        self.driver.get(self.oceanurl)
        loginfield = WebDriverWait(self.driver, 10).until(
            EC.presence_of_element_located((By.ID, "login-field"))
        )
        loginfield.send_keys('station')
        self.driver.find_element(By.ID, "password-field").send_keys('viewer')
        self.driver.find_element(By.ID, "login-btn").click()
        element = WebDriverWait(self.driver, 10).until(
            EC.presence_of_element_located((By.TAG_NAME, "svg"))
        )

    def end_pic(self):
        self.driver.close()
        pass

    def take_pic(self, url, width=None, height=None):
        if width is None or height is None:
            self.driver.maximize_window()
        else:
            self.driver.set_window_size(width, height + 134)

        self.driver.get(url)
        sleep(self.wait)
        ret = self.driver.get_screenshot_as_base64()
        return io.BytesIO(base64.urlsafe_b64decode(ret))


def generate_pptx(filename, query, chrome, **args):
    cypher_path = CypherPath()
    nodes = cypher_path.run(query, **args)
    print('length', len(nodes))

    prs = Presentation()

    prev_url = ''
    # for p in [nodes[-1]]:
    for i, p in enumerate(nodes):
        main = p['mainPath']
        interaction = p.get('interactions', {})
        url = interaction.get('url', None) if interaction else None
        href = None
        if url is None: url = main.get('url', None)
        if url is not None:
            url = eval(url)
            href = url['href']
        if main.has_label('H_UPDATE') and main['label'] == 'intention':
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = 'Starting inquiry'
            slide.placeholders[1].text = main['text']
        # elif 'url' in p and (p.has_label('H_UPDATE') or ('event' in p and p['url'] != prev_url)):
        elif href is not None and (main.has_label('H_UPDATE') or href != prev_url):
            prev_url = href
            text = interaction['label'] if interaction else main.get('event', main['label'])
            pic = chrome.take_pic(url['href'], url['screenWidth'], url['screenHeight']) if chrome is not None else None

            picture_slide_layout = prs.slide_layouts[8]
            slide = prs.slides.add_slide(picture_slide_layout)
            picture_placeholder = slide.placeholders[1]
            shapes = slide.shapes
            title = slide.placeholders[0]
            s_text = slide.placeholders[2]
            title.text = ''
            if main.has_label('H_UPDATE'):
                if i == len(nodes) - 1: title.text += 'Final Insight: '
                else: title.text += 'Insight: '
                # title.text += main.get('text', text)
                s_text.text += main.get('text', text)
            elif interaction is not None and interaction.has_label('C_UPDATE'):
                title.text += 'Interaction: '
                # title.text += main.get('text', text)
                s_text.text += main.get('text', text)
            else:
                title.text = text
            title.left, s_text.left = Inches(1), Inches(1)
            title.width, s_text.width = prs.slide_width - Inches(2), prs.slide_width - Inches(2)
            title.height, s_text.height = Inches(0.5), Inches(1)
            title.top, s_text.top = prs.slide_height - Inches(3), prs.slide_height - Inches(2)

            # slide.placeholders[0].text = 'Final Insight: ' + p['text'] if p.has_label('H_UPDATE') else p['event']
            # slide.placeholders[0].text = text
            # slide.placeholders[1].text = title.text
            # slide.placeholders[2].text = title.text

            deltax = prs.slide_width / int(url['screenWidth'])
            deltay = deltax
            if chrome is not None:
                # pic_shape = shapes.add_picture(pic, 0, 0, prs.slide_width, url['screenHeight'] * deltay) if chrome is not None else None
                picture_placeholder.left, picture_placeholder.top, picture_placeholder.width, picture_placeholder.height = 0, 0, prs.slide_width, int(url['screenHeight'] * deltay)
                pic_shape = picture_placeholder.insert_picture(pic)
                pic_shape.left, pic_shape.top, pic_shape.width, pic_shape.height = 0, 0, prs.slide_width, int(url['screenHeight'] * deltay)

            title.top, s_text.top = int(url['screenHeight'] * deltay) + Inches(0.3), int(url['screenHeight'] * deltay) + Inches(0.85)
            print(deltay)
            for s in eval(main.get('shapes', '[]')):
                if s['type'] == 'circle':
                    shape = shapes.add_shape(
                        MSO_AUTO_SHAPE_TYPE.OVAL,
                        (s['x'] - s['rx']) * deltax,
                        (s['y'] - s['ry']) * deltay,
                        (s['rx'] * 2) * deltax,
                        (s['ry'] * 2) * deltay
                    )
                else:
                    x1, x2, y1, y2 = s['x1'] * deltax, s['x2'] * deltax, (s['y1']) * deltay, (s['y2']) * deltay
                    if not takepictures:
                        shape2 = shapes.add_shape(
                            MSO_AUTO_SHAPE_TYPE.OVAL,
                            x1 - Inches(1)/2,
                            y1 - Inches(1)/2,
                            Inches(1),
                            Inches(1)
                        )
                        shape1 = shapes.add_shape(
                            MSO_AUTO_SHAPE_TYPE.OVAL,
                            x2 - Inches(1)/2,
                            y2 - Inches(1)/2,
                            Inches(1),
                            Inches(1)
                        )
                        shape1.fill.background()
                        shape2.fill.background()

                    w = Inches(0.2)
                    top = y1
                    left = x1
                    dist = math.hypot(x2 - x1, y2 - y1) + (deltax * 10)
                    ang = -(math.atan2(y2 - y1, x2 - x1))
                    deltatop = math.sin(ang) * (dist - w * math.cos(ang)) / 2
                    deltaleft = math.tan((math.pi * 2 - ang) / 2) * (deltatop - w * math.sin(ang) / 2)
                    print(top, deltatop, math.degrees(ang))
                    shape = shapes.add_shape(
                        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, left + deltaleft, top - deltatop, dist, w
                    )
                    shape.rotation = -math.degrees(ang)
                shape.fill.background()
    # if 'ref' not in e_node:
    #     blank_slide_layout = prs.slide_layouts[1]
    #     slide = prs.slides.add_slide(blank_slide_layout)
    #     slide.shapes.title.text = 'Final Insight'
    #     slide.placeholders[1].text = e_node['text']
    # else:
    #     blank_slide_layout = prs.slide_layouts[8]
    #     slide = prs.slides.add_slide(blank_slide_layout)
    #     slide.shapes.title.text = 'Final Insight: ' + e_node['text']
    #     # slide.placeholders[1].text = e_node['text']
    #     pic = slide.shapes.add_picture('local_storage/images/'+e_node['ref']+'.jpg', 0, 0, height=Inches(5))

    prs.save(filename)
    print('saved')


if __name__ == '__main__':
    takepictures = True
    chrome = None
    if takepictures:
        chrome = Chrome(15)
        chrome.setup_pic_wbt()
        # chrome.setup_pic_ocean()

    # examples of ppt generation:
    generate_pptx('../local_storage/test.pptx',
                  # CypherPath.INSIGHT_FROM_CLOSEST_NEAR_INTENTION,
                  # CypherPath.ALL_INSIGHTS_FROM_INTENTION,
                  CypherPath.INSIGHT_INTENTION,
                  # CypherPath.ONLY_INSIGHT,
                  # CypherPath.INSIGHT_FROM_LONGEST_NEAR_INTENTION,
                  # insight='B0C lack opportunities to take formal education',
                  # insight='B0C (Northern CBRM) seems to have the lowest satisfaction with access to ed',
                  # insight='region "BOC" has a very low perception of access to educational opportuniti',
                  # insight='south seems to have more sensors',
                  intention='status of health in nova scotia',
                  insight='B0C seems to be better related to its health perception',
                  # intention='what area has concerns regarding access to education',
                  chrome=chrome
                  )

    if takepictures:
        chrome.end_pic()
    print('END!')
