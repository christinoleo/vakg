#%%
import base64
import math
from time import sleep
from pptx import Presentation
from pptx.util import Inches
from pptx.shapes.autoshape import Shape
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

from app.graph.connector2 import NeoSession
from credentials import *

neodbsession = NeoSession(url, login, password)
db = neodbsession.engine.begin()
# db.run('MATCH (n) RETURN *')

from pptx import Presentation

#%%
# with neodbsession.get_db() as db:
# neodbsession.engine.rollback(db)
# db = neodbsession.engine.begin()
insighttext = 'expensive!'
# path = db.run(f'''
# match p=((n:H_UPDATE {{text:'{insighttext}'}})<-[:DOES]-()<-[:INSIGHT*1..20]-()-[:LEADS_TO]-(n1:H_UPDATE {{label:'intention'}}))
# match p1=(
#     (n)-[:DOES]-()-[:FEEDBACK]-(m:C_STATE)
#     -[:UPDATE*0..20]-(c)-[:FEEDBACK]-()-[:DOES]-(n1)
#     )
# return p1
# ''').data()
def get_path(start):
    db = neodbsession.engine.begin()
    path = db.run(f'''
    match p=((n:H_UPDATE {{text:'{start}'}})<-[:DOES]-()<-[:INSIGHT*0..20]-()-[:LEADS_TO]-(n1:H_UPDATE {{label:'intention'}}))
    match p1=(
        (n)-[:LEADS_TO]-()-[:FEEDBACK]-(m:C_STATE)
        -[:UPDATE*0..20]-(c)-[:FEEDBACK]-()-[:LEADS_TO]-(n1)
        )
    return p1
    ''').data()
    print(path)
    nodes = path[0]['p1'].nodes[::-1]
    neodbsession.engine.rollback(db)
    return nodes
# path = get_path(insighttext)

# db.run('match (v) return v')
# match (n:H_UPDATE {label:'intention'}),(r) where not (n)<--(r) return n
# match (n:H_UPDATE {label:'insight'}) return n.text
# match (n:H_UPDATE {label:'intention'}) return n.text

# path to closest intention
# match p=((n:H_UPDATE {text:'close to a third of NS complain of education prices'})<-[:FOLLOWS_INSIGHT*1..20]-(n1:H_UPDATE {label:'intention'})) return p limit 1

# path to closest interaction
# match p=((n:H_UPDATE {text:'close to a third of NS complain of education prices'})-[:DOES]-()-[:FEEDBACK]-()-[:LEADS_TO]-(m:C_UPDATE {user:n.user})) return *

# interaction path between insight and intention
# match p=((n:H_UPDATE {text:'close to a third of NS complain of education prices'})<-[:FOLLOWS_INSIGHT*1..20]-(n1:H_UPDATE {label:'intention'}))
# match p1=(
#     (n)-[:DOES]-()-[:FEEDBACK]-()-[:LEADS_TO]-(m:C_UPDATE {user:n.user})
#     -[:FOLLOWS_UPDATE*1..20]-(c)-[:LEADS_TO]-()-[:FEEDBACK]-()-[:DOES]-(n1)
#     )
# return p1


#%%
from selenium import webdriver
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.common.by import By
from selenium.webdriver.support.wait import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC

urls = [
"http://localhost:3000/?tree=0a0b1a0b1b5a0b1b5b0b0&map=EDOPP_C1&maximized=0",
"http://localhost:3000/?tree=0a0b1a0b1b5a0b1b5b0b0&map=EDOPP_C1&maximized=1"
]

driver = webdriver.Chrome('../local_storage/chromedriver.exe')
driver.maximize_window()
driver.get('http://localhost:3000/')
element = WebDriverWait(driver, 10).until(
        EC.presence_of_element_located((By.ID, 'disable-tutorial'))
    )
element.click()
element = WebDriverWait(driver, 10).until(
        EC.presence_of_element_located((By.TAG_NAME, "path"))
    )

for u in urls:
    driver.get(u)
    sleep(3)
    a = driver.get_screenshot_as_png()
# with open(f"local_storage/images/ss.png", "wb") as fh:
#     fh.write(base64.urlsafe_b64decode(driver.get_screenshot_as_base64()))
driver.close()
#%%

img_path = 'local_storage/images/ss.png'

prs = Presentation()
blank_slide_layout = prs.slide_layouts[8]
slide = prs.slides.add_slide(blank_slide_layout)
slide.shapes.title.text = 'awliuh'
slide.placeholders[1].text = ''

# left = top = Inches(0)
# pic = slide.shapes.add_picture(img_path, 0, 0, prs.slide_width)

shapes = slide.shapes
x1 = y1 = 0
x2 = Inches(2)
y2 = Inches(4)
w = Inches(0.2)
top = min(y1, y2)
left = min(x1, x2)
dist = math.hypot(x1+x2, y1+y2)
ang = math.radians(360) - (math.radians(270)+math.atan2(y2-y1, x2-x1))
left = left+math.sin(ang)*(dist-w*math.cos(ang))/2
top = top+math.tan((math.pi*2-ang)/2)*(left+w*math.sin(ang)/2)
shape = shapes.add_shape(
    MSO_AUTO_SHAPE_TYPE.DOWN_ARROW, left, top, w, dist
)
shape.rotation = 360 - math.degrees(ang)
shape.fill.background()

prs.save('local_storage/test.pptx')
#%%
import io

class Chrome:
    def __init__(self):
        self.driver = webdriver.Chrome('../local_storage/chromedriver.exe')
        self.url = 'http://localhost:3000/'

    def setup_pic(self):
        self.driver.get(self.url)
        element = WebDriverWait(self.driver, 10).until(
            EC.presence_of_element_located((By.ID, 'disable-tutorial'))
        )
        element.click()
        element = WebDriverWait(self.driver, 10).until(
            EC.presence_of_element_located((By.TAG_NAME, "path"))
        )

    def end_pic(self):
        self.driver.close()

    def take_pic(self, url, width=None, height=None):
        if width is None or height is None:
            self.driver.maximize_window()
        else:
            self.driver.set_window_size(width, height)

        self.driver.get(url)
        sleep(3)
        ret = self.driver.get_screenshot_as_base64()
        return io.BytesIO(base64.urlsafe_b64decode(ret))

nodes = get_path('expensive!')
s_node = nodes[0]
e_node = nodes[-1]

prs = Presentation()
blank_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(blank_slide_layout)
slide.shapes.title.text = 'Starting inquiry'
slide.placeholders[1].text = s_node['text']

chrome = Chrome()
chrome.setup_pic()
prev_url = ''
# for p in [nodes[-1]]:
for p in nodes[1:]:
    print(p)
    if 'url' in p and (p.has_label('H_UPDATE') or ('event' in p and p['url'] != prev_url)):
        pic = chrome.take_pic(p['url'], p.get('screenWidth', None), p.get('screenHeight', None))

        blank_slide_layout = prs.slide_layouts[8]
        slide = prs.slides.add_slide(blank_slide_layout)
        shapes = slide.shapes
        title = shapes.title
        title.text = 'Final Insight: ' + p['text'] if p.has_label('H_UPDATE') else p['event']
        # slide.placeholders[0].text = 'Final Insight: ' + p['text'] if p.has_label('H_UPDATE') else p['event']
        slide.placeholders[0].text = title.text
        slide.placeholders[1].text = title.text
        slide.placeholders[2].text = title.text
        pic_shape = shapes.add_picture(pic, 0, 0, prs.slide_width)
        deltax = pic_shape.width/int(p.get('screenWidth', 1))
        deltay = pic_shape.height/int(p.get('screenHeight', 1))
        for s in eval(p.get('shapes', '[]')):
            if s['type'] == 'circle':
                shape = shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.OVAL,
                    (s['x']-s['rx'])*deltax,
                    (s['y']-s['ry'])*deltay,
                    (s['rx']*2)*deltax,
                    (s['ry']*2)*deltay
                )
            else:
                x1,x2,y1,y2 = s['x1']*deltax,s['x2']*deltax,s['y1']*deltay,s['y2']*deltay
                w = Inches(0.2)
                top = min(y1, y2)
                left = min(x1, x2)
                dist = math.hypot(x2-x1, y2-y1)
                ang = math.radians(360) - (math.radians(270) + math.atan2(y2 - y1, x2 - x1))
                deltaleft = math.sin(ang) * (dist - w * math.cos(ang)) / 2
                left = left + deltaleft
                # top = top + math.tan((math.pi * 2 - ang) / 2) * (deltaleft + w * math.sin(ang) / 2)
                shape = shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.DOWN_ARROW, left, top, w, dist
                )
                shape.rotation = 360 - math.degrees(ang)
            shape.fill.background()
        prev_url = p['url']
# if 'ref' not in e_node:
#     blank_slide_layout = prs.slide_layouts[1]
#     slide = prs.slides.add_slide(blank_slide_layout)
#     slide.shapes.title.text = 'Final Insight'
#     slide.placeholders[1].text = e_node['text']
# else:
#     blank_slide_layout = prs.slide_layouts[8]
#     slide = prs.slides.add_slide(blank_slide_layout)
#     slide.shapes.title.text = 'Final Insight: ' + e_node['text']
#     # slide.placeholders[1].text = e_node['text']
#     pic = slide.shapes.add_picture('local_storage/images/'+e_node['ref']+'.jpg', 0, 0, height=Inches(5))

chrome.end_pic()
prs.save('local_storage/test.pptx')